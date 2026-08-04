"""MindSurve-branded PowerPoint export for study analytics."""

from __future__ import annotations

import io
import re
from datetime import datetime, timezone
from typing import Any, Dict, List, Optional, Sequence, Tuple

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.util import Inches, Pt
from sqlalchemy.orm import Session

from app.models.study_model import Study
from app.schemas.assistant_schema import (
    AssistantMetric,
    AssistantQueryPlan,
    AppliedContext,
    RankDirection,
)
from app.services.assistant_tools import (
    _biggest_age_element_gap,
    _biggest_gender_element_gap,
    _classification_questions,
    metric_prefix,
    tool_classification_distribution,
    tool_executive_summary,
    tool_fatigue_summary,
    tool_rank_designs,
    tool_rank_elements,
    tool_response_time_summary,
    tool_use_avoid,
)
from app.services.ppt_images import (
    ImageCache,
    compose_design_preview,
    compose_element_thumb,
)

# Website theme (#2674BA family)
BRAND_BLUE = RGBColor(0x26, 0x74, 0xBA)
BRAND_BLUE_DARK = RGBColor(0x1F, 0x5F, 0x99)
BRAND_BLUE_SOFT = RGBColor(0xE8, 0xF1, 0xF9)
BRAND_BLUE_MID = RGBColor(0xA8, 0xC8, 0xE8)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
NEAR_BLACK = RGBColor(0x11, 0x18, 0x27)
SLATE = RGBColor(0x4B, 0x55, 0x63)
MUTED = RGBColor(0x6B, 0x72, 0x80)
LIGHT_LINE = RGBColor(0xE5, 0xE7, 0xEB)
ACCENT_GREEN = RGBColor(0x05, 0x96, 0x69)
ACCENT_RED = RGBColor(0xDC, 0x26, 0x26)
SLIDE_W = Inches(13.333)
SLIDE_H = Inches(7.5)


def _set_run_font(run, *, size: int, bold: bool = False, color: RGBColor = NEAR_BLACK, name: str = "Calibri") -> None:
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.color.rgb = color
    run.font.name = name


def _add_rect(slide, left, top, width, height, fill: RGBColor, line: Optional[RGBColor] = None):
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill
    if line is None:
        shape.line.fill.background()
    else:
        shape.line.color.rgb = line
    return shape


def _add_round_rect(slide, left, top, width, height, fill: RGBColor):
    shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill
    shape.line.fill.background()
    try:
        shape.adjustments[0] = 0.08
    except Exception:
        pass
    return shape


def _textbox(slide, left, top, width, height, text: str, *, size: int = 14, bold: bool = False,
             color: RGBColor = NEAR_BLACK, align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.TOP):
    box = slide.shapes.add_textbox(left, top, width, height)
    tf = box.text_frame
    tf.word_wrap = True
    tf.auto_size = None
    try:
        tf._txBody.bodyPr.set("anchor", {MSO_ANCHOR.TOP: "t", MSO_ANCHOR.MIDDLE: "ctr", MSO_ANCHOR.BOTTOM: "b"}.get(anchor, "t"))
    except Exception:
        pass
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text or ""
    _set_run_font(run, size=size, bold=bold, color=color)
    return box


def _multi_lines(slide, left, top, width, height, lines: Sequence[Tuple[str, dict]]):
    box = slide.shapes.add_textbox(left, top, width, height)
    tf = box.text_frame
    tf.word_wrap = True
    for idx, (text, style) in enumerate(lines):
        p = tf.paragraphs[0] if idx == 0 else tf.add_paragraph()
        p.alignment = style.get("align", PP_ALIGN.LEFT)
        p.space_after = Pt(style.get("space_after", 6))
        run = p.add_run()
        run.text = text or ""
        _set_run_font(
            run,
            size=style.get("size", 14),
            bold=style.get("bold", False),
            color=style.get("color", NEAR_BLACK),
        )
    return box


def _brand_wordmark(slide, left, top, *, size: int = 28):
    """Render Mind + Surve wordmark matching the navbar (no logo image asset)."""
    box = slide.shapes.add_textbox(left, top, Inches(4.5), Inches(0.6))
    tf = box.text_frame
    p = tf.paragraphs[0]
    r1 = p.add_run()
    r1.text = "Mind"
    _set_run_font(r1, size=size, bold=True, color=BRAND_BLUE)
    r2 = p.add_run()
    r2.text = "Surve"
    _set_run_font(r2, size=size, bold=True, color=SLATE)
    return box


def _slide_chrome(slide, title: str, subtitle: Optional[str] = None, *, page: int = 0, total: int = 0):
    _add_rect(slide, 0, 0, SLIDE_W, Inches(0.08), BRAND_BLUE)
    _add_rect(slide, 0, SLIDE_H - Inches(0.42), SLIDE_W, Inches(0.42), BRAND_BLUE_SOFT)
    _brand_wordmark(slide, Inches(0.55), Inches(0.22), size=16)
    _textbox(slide, Inches(0.55), Inches(0.7), Inches(11.5), Inches(0.45), title,
             size=26, bold=True, color=NEAR_BLACK)
    if subtitle:
        _textbox(slide, Inches(0.55), Inches(1.15), Inches(11.5), Inches(0.35), subtitle,
                 size=13, color=MUTED)
    footer = "MindSurve AI · Analytics presentation"
    if page and total:
        footer = f"{footer}  ·  {page} / {total}"
    _textbox(
        slide,
        Inches(0.55),
        SLIDE_H - Inches(0.36),
        Inches(12),
        Inches(0.28),
        footer,
        size=10,
        color=SLATE,
    )


def _blank(prs: Presentation):
    return prs.slides.add_slide(prs.slide_layouts[6])


def _safe(text: Any, limit: int = 120) -> str:
    value = " ".join(str(text or "").split())
    if len(value) <= limit:
        return value
    return value[: max(0, limit - 1)].rstrip() + "…"


def _fmt_num(value: Any, digits: int = 2) -> str:
    try:
        num = float(value)
    except (TypeError, ValueError):
        return str(value or "—")
    if abs(num - round(num)) < 1e-9:
        return str(int(round(num)))
    return f"{num:.{digits}f}"


def _items_from_rank_result(result: Dict[str, Any]) -> List[Dict[str, Any]]:
    for block in result.get("blocks") or []:
        data = block.get("data") or {}
        items = data.get("items")
        if isinstance(items, list) and items:
            return items
        designs = data.get("designs")
        if isinstance(designs, list) and designs:
            return designs
    return []


def _design_pack_from_rank_result(result: Dict[str, Any]) -> Dict[str, Any]:
    """Extract designs plus background/aspect/study_type from a rank_designs result."""
    for block in result.get("blocks") or []:
        data = block.get("data") or {}
        designs = data.get("designs") or data.get("items") or []
        if isinstance(designs, list) and designs:
            return {
                "designs": designs,
                "background_url": data.get("background_url"),
                "aspect_ratio": data.get("aspect_ratio") or "9 / 16",
                "study_type": data.get("study_type") or "grid",
            }
    return {"designs": [], "background_url": None, "aspect_ratio": "9 / 16", "study_type": "grid"}


def _add_picture_bytes(slide, png_bytes: bytes, left, top, width=None, height=None):
    stream = io.BytesIO(png_bytes)
    kwargs = {}
    if width is not None:
        kwargs["width"] = width
    if height is not None:
        kwargs["height"] = height
    return slide.shapes.add_picture(stream, left, top, **kwargs)


def _classification_rows(db: Session, study_obj: Study, plan: AssistantQueryPlan) -> List[Dict[str, Any]]:
    rows: List[Dict[str, Any]] = []
    for question in _classification_questions(study_obj)[:4]:
        q_plan = plan.model_copy(
            update={
                "classification_question": question.question_text,
                "classification_options": [],
            }
        )
        result = tool_classification_distribution(db, study_obj, q_plan, None)
        if result.get("status") == "needs_clarification":
            continue
        for block in result.get("blocks") or []:
            if block.get("type") != "classification_distribution":
                continue
            data = block.get("data") or {}
            options = data.get("options") or []
            rows.append(
                {
                    "question": data.get("question") or question.question_text,
                    "answered": data.get("answered") or 0,
                    "total": data.get("total_respondents") or 0,
                    "options": options[:8],
                }
            )
    return rows


def _kpi_card(slide, left, top, width, height, label: str, value: str, accent: RGBColor = BRAND_BLUE):
    _add_round_rect(slide, left, top, width, height, WHITE)
    _add_rect(slide, left, top, Inches(0.08), height, accent)
    _textbox(slide, left + Inches(0.22), top + Inches(0.18), width - Inches(0.3), Inches(0.3),
             label, size=11, color=MUTED)
    _textbox(slide, left + Inches(0.22), top + Inches(0.5), width - Inches(0.3), Inches(0.55),
             value, size=26, bold=True, color=NEAR_BLACK)


def build_analytics_pptx(
    *,
    db: Session,
    study_obj: Study,
    analysis: Dict[str, Any],
    plan: AssistantQueryPlan,
    context: AppliedContext,
) -> Tuple[bytes, Dict[str, Any]]:
    """Build a ~15-slide MindSurve analytics deck. Returns (bytes, meta)."""
    prs = Presentation()
    prs.slide_width = SLIDE_W
    prs.slide_height = SLIDE_H

    study_title = context.study_title or study_obj.title or "Study"
    study_type = str(study_obj.study_type or context.study_type or "grid").title()
    metric = plan.metric or AssistantMetric.T
    summary = analysis.get("dashboard_summary") or {}
    panelists = int(summary.get("uniquePanelists") or 0)
    responses = int(summary.get("totalResponses") or 0)
    avg_rating = round(float(summary.get("avgRating") or 0), 2)
    avg_rt = round(float(summary.get("avgResponseTime") or 0), 2)
    generated_at = datetime.now(timezone.utc).strftime("%d %b %Y")

    base_plan = plan.model_copy(update={"metric": metric, "limit": 5})
    image_cache = ImageCache()
    info = analysis.get("Information Block") or {}
    fallback_bg = info.get("Study Background") or info.get("background_image_url")
    fallback_aspect = info.get("Aspect Ratio") or "9 / 16"
    study_type_key = str(study_obj.study_type or context.study_type or "grid").lower()

    top_elements = _items_from_rank_result(
        tool_rank_elements(
            analysis,
            study_obj,
            base_plan.model_copy(update={"direction": RankDirection.highest, "limit": 5}),
        )
    )
    worst_elements = _items_from_rank_result(
        tool_rank_elements(
            analysis,
            study_obj,
            base_plan.model_copy(update={"direction": RankDirection.lowest, "limit": 5}),
        )
    )
    top_design_pack = _design_pack_from_rank_result(
        tool_rank_designs(
            analysis,
            study_obj,
            base_plan.model_copy(update={"direction": RankDirection.highest, "limit": 5, "must_include": []}),
        )
    )
    worst_design_pack = _design_pack_from_rank_result(
        tool_rank_designs(
            analysis,
            study_obj,
            base_plan.model_copy(update={"direction": RankDirection.lowest, "limit": 5, "must_include": []}),
        )
    )
    exec_result = tool_executive_summary(db, analysis, study_obj, base_plan, context)
    exec_bullets = []
    for block in exec_result.get("blocks") or []:
        if block.get("type") == "executive_summary":
            exec_bullets = (block.get("data") or {}).get("bullets") or []
    use_result = tool_use_avoid(analysis, study_obj, base_plan)
    use_block = next((b for b in use_result.get("blocks") or [] if b.get("type") == "use_avoid"), None)
    use_items = ((use_block or {}).get("data") or {}).get("use") or []
    avoid_items = ((use_block or {}).get("data") or {}).get("avoid") or []
    class_rows = _classification_rows(db, study_obj, base_plan)
    metric_code = metric_prefix(metric.value)
    gender_gap = _biggest_gender_element_gap(analysis, metric_code)
    age_gap = _biggest_age_element_gap(analysis, metric_code)
    rt_result = tool_response_time_summary(analysis)
    fatigue_result = tool_fatigue_summary(analysis)
    rating_dist = summary.get("ratingDistribution") or []

    def cover():
        slide = _blank(prs)
        _add_rect(slide, 0, 0, SLIDE_W, SLIDE_H, WHITE)
        _add_rect(slide, 0, 0, Inches(0.35), SLIDE_H, BRAND_BLUE)
        _add_rect(slide, Inches(0.35), 0, Inches(0.08), SLIDE_H, BRAND_BLUE_MID)
        _add_rect(slide, 0, SLIDE_H - Inches(1.1), SLIDE_W, Inches(1.1), BRAND_BLUE)
        _brand_wordmark(slide, Inches(1.0), Inches(1.6), size=40)
        _textbox(slide, Inches(1.0), Inches(2.5), Inches(10.5), Inches(0.4),
                 "Analytics presentation", size=16, color=MUTED)
        _textbox(slide, Inches(1.0), Inches(3.2), Inches(11), Inches(1.2),
                 study_title, size=36, bold=True, color=NEAR_BLACK)
        meta = f"{study_type} study  ·  {panelists} panelists  ·  {responses} responses  ·  {generated_at}"
        _textbox(slide, Inches(1.0), Inches(4.7), Inches(11), Inches(0.4), meta, size=13, color=SLATE)
        _textbox(
            slide,
            Inches(1.0),
            SLIDE_H - Inches(0.75),
            Inches(11),
            Inches(0.35),
            "Verified insights from MindSurve AI",
            size=12,
            color=WHITE,
        )

    def agenda():
        slide = _blank(prs)
        _slide_chrome(slide, "Agenda", "What this deck covers", page=2, total=15)
        items = [
            "1. Study overview & respondent profile",
            "2. Key performance indicators",
            "3. Executive findings",
            "4. Strongest & weakest elements",
            "5. Top & bottom designs",
            "6. Classification insights",
            "7. Segment differences",
            "8. Use / avoid guidance",
            "9. Response behaviour & next steps",
        ]
        y = Inches(1.7)
        for item in items:
            _add_round_rect(slide, Inches(0.7), y, Inches(11.8), Inches(0.48), WHITE)
            _add_rect(slide, Inches(0.7), y, Inches(0.08), Inches(0.48), BRAND_BLUE)
            _textbox(slide, Inches(1.0), y + Inches(0.08), Inches(11), Inches(0.35), item,
                     size=15, color=NEAR_BLACK)
            y += Inches(0.54)

    def overview():
        slide = _blank(prs)
        _slide_chrome(slide, "Study overview", study_title, page=3, total=15)
        cards = [
            ("Panelists", str(panelists)),
            ("Responses", str(responses)),
            ("Avg rating", _fmt_num(avg_rating)),
            ("Avg response time", f"{_fmt_num(avg_rt)}s"),
        ]
        x = Inches(0.7)
        for label, value in cards:
            _kpi_card(slide, x, Inches(1.8), Inches(2.85), Inches(1.35), label, value)
            x += Inches(3.05)
        _add_round_rect(slide, Inches(0.7), Inches(3.5), Inches(11.8), Inches(2.8), WHITE)
        _textbox(slide, Inches(1.0), Inches(3.7), Inches(11), Inches(0.35),
                 "Study snapshot", size=16, bold=True, color=BRAND_BLUE)
        lines = [
            (f"Study type: {study_type}", {"size": 14, "color": NEAR_BLACK, "space_after": 8}),
            (f"Metric focus: Top Down (T)", {"size": 14, "color": NEAR_BLACK, "space_after": 8}),
            (f"Segment: {context.segment_label or 'Overall'}", {"size": 14, "color": NEAR_BLACK, "space_after": 8}),
            (
                "All figures in this deck are computed from verified study analysis — not model estimates.",
                {"size": 13, "color": MUTED, "space_after": 0},
            ),
        ]
        _multi_lines(slide, Inches(1.0), Inches(4.2), Inches(11), Inches(1.8), lines)

    def respondent_profile():
        slide = _blank(prs)
        _slide_chrome(slide, "Respondent profile", "Who answered this study", page=4, total=15)
        _kpi_card(slide, Inches(0.7), Inches(1.8), Inches(3.7), Inches(1.4), "Unique panelists", str(panelists))
        _kpi_card(slide, Inches(4.7), Inches(1.8), Inches(3.7), Inches(1.4), "Scored responses", str(responses), BRAND_BLUE_DARK)
        _kpi_card(slide, Inches(8.7), Inches(1.8), Inches(3.7), Inches(1.4), "Avg rating", _fmt_num(avg_rating), ACCENT_GREEN)
        _add_round_rect(slide, Inches(0.7), Inches(3.5), Inches(11.8), Inches(2.8), WHITE)
        _textbox(slide, Inches(1.0), Inches(3.7), Inches(11), Inches(0.35),
                 "Rating distribution", size=16, bold=True, color=BRAND_BLUE)
        if rating_dist:
            lines = []
            for item in rating_dist[:8]:
                name = _safe(item.get("name") or item.get("label"), 40)
                value = item.get("value")
                lines.append((f"• {name}: {_fmt_num(value)}", {"size": 14, "space_after": 6}))
            _multi_lines(slide, Inches(1.0), Inches(4.2), Inches(11), Inches(1.8), lines)
        else:
            _textbox(slide, Inches(1.0), Inches(4.4), Inches(11), Inches(0.4),
                     "Rating distribution is not available for this cohort yet.", size=13, color=MUTED)

    def executive():
        slide = _blank(prs)
        _slide_chrome(slide, "Executive findings", "Most important verified takeaways", page=5, total=15)
        if not exec_bullets:
            _textbox(slide, Inches(0.7), Inches(2.0), Inches(11.5), Inches(1),
                     "Not enough verified analytics are available yet for an executive summary.",
                     size=14, color=MUTED)
            return
        y = Inches(1.7)
        for bullet in exec_bullets[:5]:
            _add_round_rect(slide, Inches(0.7), y, Inches(11.8), Inches(0.95), WHITE)
            _add_rect(slide, Inches(0.7), y, Inches(0.1), Inches(0.95), BRAND_BLUE)
            title = _safe(bullet.get("title") or f"Finding {bullet.get('rank')}", 40)
            text = _safe(bullet.get("text"), 140)
            text_left = Inches(1.05)
            if bullet.get("image_url"):
                try:
                    thumb = compose_element_thumb(
                        {"name": title, "image_url": bullet.get("image_url")},
                        cache=image_cache,
                        size=160,
                    )
                    _add_picture_bytes(
                        slide,
                        thumb,
                        Inches(0.95),
                        y + Inches(0.14),
                        width=Inches(0.68),
                        height=Inches(0.68),
                    )
                    text_left = Inches(1.85)
                except Exception:
                    pass
            _textbox(slide, text_left, y + Inches(0.12), Inches(10.2), Inches(0.3),
                     title, size=13, bold=True, color=BRAND_BLUE)
            _textbox(slide, text_left, y + Inches(0.45), Inches(10.2), Inches(0.4),
                     text, size=13, color=NEAR_BLACK)
            y += Inches(1.02)

    def elements_slide(title: str, subtitle: str, items: List[Dict[str, Any]], *, page: int, accent: RGBColor):
        """One slide with all elements as image + name + score cards."""
        slide = _blank(prs)
        _slide_chrome(slide, title, subtitle, page=page, total=15)
        cards = list(items or [])[:5]
        if not cards:
            _textbox(slide, Inches(0.7), Inches(2.0), Inches(11.5), Inches(0.5),
                     "No element rankings are available yet.", size=14, color=MUTED)
            return

        count = len(cards)
        gap = Inches(0.22)
        left0 = Inches(0.55)
        usable = Inches(12.2)
        card_w = (usable - gap * (count - 1)) / count
        card_top = Inches(1.7)
        card_h = Inches(5.0)
        img_pad = Inches(0.14)
        img_h = Inches(3.15)

        for idx, item in enumerate(cards):
            left = left0 + (card_w + gap) * idx
            _add_round_rect(slide, left, card_top, card_w, card_h, WHITE)
            _add_rect(slide, left, card_top, card_w, Inches(0.08), accent)
            try:
                thumb = compose_element_thumb(item, cache=image_cache, size=420)
                _add_picture_bytes(
                    slide,
                    thumb,
                    left + img_pad,
                    card_top + Inches(0.28),
                    width=card_w - img_pad * 2,
                    height=img_h,
                )
            except Exception:
                _textbox(
                    slide,
                    left + img_pad,
                    card_top + Inches(1.2),
                    card_w - img_pad * 2,
                    Inches(0.5),
                    "No image",
                    size=12,
                    color=MUTED,
                    align=PP_ALIGN.CENTER,
                )

            rank = item.get("rank") or idx + 1
            name = _safe(item.get("name"), 36)
            category = _safe(item.get("category") or item.get("category_name"), 28)
            value = _fmt_num(item.get("value"))
            text_top = card_top + Inches(3.55)
            _textbox(
                slide,
                left + Inches(0.1),
                text_top,
                card_w - Inches(0.2),
                Inches(0.28),
                f"#{rank}  ·  {value}",
                size=12,
                bold=True,
                color=accent,
                align=PP_ALIGN.CENTER,
            )
            _textbox(
                slide,
                left + Inches(0.1),
                text_top + Inches(0.32),
                card_w - Inches(0.2),
                Inches(0.7),
                name,
                size=12,
                bold=True,
                color=NEAR_BLACK,
                align=PP_ALIGN.CENTER,
            )
            if category:
                _textbox(
                    slide,
                    left + Inches(0.1),
                    text_top + Inches(1.05),
                    card_w - Inches(0.2),
                    Inches(0.3),
                    category,
                    size=10,
                    color=MUTED,
                    align=PP_ALIGN.CENTER,
                )

    def designs_slide(
        title: str,
        subtitle: str,
        pack: Dict[str, Any],
        *,
        page: int,
        accent: RGBColor,
    ):
        """One slide with all designs as configurator-style previews + score."""
        slide = _blank(prs)
        _slide_chrome(slide, title, subtitle, page=page, total=15)
        designs = list(pack.get("designs") or [])[:5]
        if not designs:
            _textbox(slide, Inches(0.7), Inches(2.0), Inches(11.5), Inches(0.5),
                     "No complete designs could be ranked for this study yet.", size=14, color=MUTED)
            return

        bg_url = pack.get("background_url") or fallback_bg
        aspect = pack.get("aspect_ratio") or fallback_aspect
        dtype = str(pack.get("study_type") or study_type_key).lower()
        count = len(designs)
        gap = Inches(0.2)
        left0 = Inches(0.5)
        usable = Inches(12.3)
        card_w = (usable - gap * (count - 1)) / count
        card_top = Inches(1.65)
        card_h = Inches(5.1)
        # Portrait designs are taller; landscape wider — fit inside card.
        preview_h = Inches(3.55) if dtype == "layer" and "16" not in str(aspect).replace(" ", "") else Inches(3.2)
        preview_top = card_top + Inches(0.22)

        for idx, design in enumerate(designs):
            left = left0 + (card_w + gap) * idx
            _add_round_rect(slide, left, card_top, card_w, card_h, WHITE)
            _add_rect(slide, left, card_top, card_w, Inches(0.08), accent)
            try:
                preview = compose_design_preview(
                    design,
                    study_type=dtype,
                    background_url=bg_url,
                    aspect_ratio=aspect,
                    cache=image_cache,
                )
                # Keep aspect: for portrait, constrain by height; for landscape by width.
                aspect_token = str(aspect or "").replace(" ", "")
                if "16/9" in aspect_token or "16:9" in aspect_token:
                    pic_w = card_w - Inches(0.22)
                    _add_picture_bytes(slide, preview, left + Inches(0.11), preview_top, width=pic_w)
                else:
                    pic_h = preview_h
                    # Center narrower portrait preview in the card.
                    # Estimate width from 9:16 ≈ height * 9/16
                    est_w = pic_h * 9 / 16
                    pic_left = left + (card_w - est_w) / 2
                    _add_picture_bytes(slide, preview, pic_left, preview_top, height=pic_h)
            except Exception:
                _textbox(
                    slide,
                    left + Inches(0.15),
                    preview_top + Inches(1.2),
                    card_w - Inches(0.3),
                    Inches(0.4),
                    "Preview unavailable",
                    size=11,
                    color=MUTED,
                    align=PP_ALIGN.CENTER,
                )

            rank = design.get("rank") or idx + 1
            score = _fmt_num(design.get("score"))
            elements = design.get("elements") or []
            names = ", ".join(_safe(el.get("name"), 18) for el in elements[:3])
            if len(elements) > 3:
                names += f" +{len(elements) - 3}"
            text_top = card_top + Inches(3.95)
            _textbox(
                slide,
                left + Inches(0.08),
                text_top,
                card_w - Inches(0.16),
                Inches(0.28),
                f"#{rank}  ·  {score}",
                size=12,
                bold=True,
                color=accent,
                align=PP_ALIGN.CENTER,
            )
            _textbox(
                slide,
                left + Inches(0.08),
                text_top + Inches(0.32),
                card_w - Inches(0.16),
                Inches(0.7),
                names or "Complete design",
                size=10,
                color=NEAR_BLACK,
                align=PP_ALIGN.CENTER,
            )

    def classification():
        slide = _blank(prs)
        _slide_chrome(slide, "Classification insights", "How respondents answered screening questions", page=10, total=15)
        if not class_rows:
            _textbox(slide, Inches(0.7), Inches(2.0), Inches(11.5), Inches(0.5),
                     "This study has no classification questions, or no answers yet.", size=14, color=MUTED)
            return
        y = Inches(1.7)
        for row in class_rows[:2]:
            _add_round_rect(slide, Inches(0.7), y, Inches(11.8), Inches(2.4), WHITE)
            _textbox(slide, Inches(1.0), y + Inches(0.15), Inches(11.2), Inches(0.35),
                     _safe(row["question"], 90), size=14, bold=True, color=BRAND_BLUE)
            meta = f"Answered {row['answered']} of {row['total']}"
            _textbox(slide, Inches(1.0), y + Inches(0.5), Inches(11.2), Inches(0.28),
                     meta, size=11, color=MUTED)
            opt_lines = []
            for opt in row["options"][:5]:
                opt_lines.append(
                    (
                        f"• {_safe(opt.get('option'), 50)} — {opt.get('count', 0)} "
                        f"({_fmt_num(opt.get('percentage'), 1)}%)",
                        {"size": 13, "space_after": 4},
                    )
                )
            if opt_lines:
                _multi_lines(slide, Inches(1.0), y + Inches(0.85), Inches(11.2), Inches(1.4), opt_lines)
            y += Inches(2.55)

    def classification_detail():
        slide = _blank(prs)
        _slide_chrome(slide, "Classification detail", "Additional screening splits", page=11, total=15)
        extra = class_rows[2:4]
        if not extra:
            # Reuse first question options in a wider list if only one/two exist.
            if class_rows:
                row = class_rows[0]
                _add_round_rect(slide, Inches(0.7), Inches(1.7), Inches(11.8), Inches(4.8), WHITE)
                _textbox(slide, Inches(1.0), Inches(1.9), Inches(11.2), Inches(0.35),
                         _safe(row["question"], 90), size=15, bold=True, color=BRAND_BLUE)
                lines = []
                for opt in row["options"][:8]:
                    lines.append(
                        (
                            f"• {_safe(opt.get('option'), 60)} — {opt.get('count', 0)} "
                            f"({_fmt_num(opt.get('percentage'), 1)}%)",
                            {"size": 14, "space_after": 8},
                        )
                    )
                _multi_lines(slide, Inches(1.0), Inches(2.5), Inches(11.2), Inches(3.6), lines or [
                    ("No option counts available.", {"size": 13, "color": MUTED})
                ])
            else:
                _textbox(slide, Inches(0.7), Inches(2.0), Inches(11.5), Inches(0.5),
                         "No further classification detail available.", size=14, color=MUTED)
            return
        y = Inches(1.7)
        for row in extra:
            _add_round_rect(slide, Inches(0.7), y, Inches(11.8), Inches(2.4), WHITE)
            _textbox(slide, Inches(1.0), y + Inches(0.15), Inches(11.2), Inches(0.35),
                     _safe(row["question"], 90), size=14, bold=True, color=BRAND_BLUE)
            lines = []
            for opt in row["options"][:5]:
                lines.append(
                    (
                        f"• {_safe(opt.get('option'), 50)} — {opt.get('count', 0)} "
                        f"({_fmt_num(opt.get('percentage'), 1)}%)",
                        {"size": 13, "space_after": 4},
                    )
                )
            _multi_lines(slide, Inches(1.0), y + Inches(0.65), Inches(11.2), Inches(1.5), lines)
            y += Inches(2.55)

    def segments():
        slide = _blank(prs)
        _slide_chrome(slide, "Segment differences", "Where audiences diverge most", page=12, total=15)
        cards: List[Tuple[str, str, str]] = []
        if gender_gap and gender_gap.get("gap", 0) > 0:
            cards.append(
                (
                    "Biggest gender gap",
                    _safe(gender_gap.get("element"), 40),
                    f"Male {_fmt_num(gender_gap.get('male'))} vs Female {_fmt_num(gender_gap.get('female'))} "
                    f"(gap {_fmt_num(gender_gap.get('gap'))})",
                )
            )
        if age_gap and age_gap.get("gap", 0) > 0:
            cards.append(
                (
                    "Biggest age gap",
                    _safe(age_gap.get("element"), 40),
                    f"{age_gap.get('left_segment')} {_fmt_num(age_gap.get('left_value'))} vs "
                    f"{age_gap.get('right_segment')} {_fmt_num(age_gap.get('right_value'))}",
                )
            )
        if not cards:
            _textbox(slide, Inches(0.7), Inches(2.0), Inches(11.5), Inches(0.6),
                     "No large gender or age element gaps were detected for this cohort.",
                     size=14, color=MUTED)
            return
        y = Inches(1.8)
        for title, element, detail in cards[:3]:
            _add_round_rect(slide, Inches(0.7), y, Inches(11.8), Inches(1.6), WHITE)
            _add_rect(slide, Inches(0.7), y, Inches(0.1), Inches(1.6), BRAND_BLUE)
            _textbox(slide, Inches(1.05), y + Inches(0.2), Inches(11), Inches(0.3),
                     title, size=12, bold=True, color=BRAND_BLUE)
            _textbox(slide, Inches(1.05), y + Inches(0.55), Inches(11), Inches(0.35),
                     element, size=18, bold=True, color=NEAR_BLACK)
            _textbox(slide, Inches(1.05), y + Inches(1.0), Inches(11), Inches(0.35),
                     detail, size=13, color=SLATE)
            y += Inches(1.85)

    def _use_avoid_column(slide, items: List[Dict[str, Any]], *, left, title: str, accent: RGBColor):
        _add_round_rect(slide, left, Inches(1.7), Inches(5.9), Inches(4.9), WHITE)
        _add_rect(slide, left, Inches(1.7), Inches(5.9), Inches(0.5), accent)
        _textbox(slide, left + Inches(0.2), Inches(1.8), Inches(5.4), Inches(0.32),
                 title, size=16, bold=True, color=WHITE)
        row = items[:5]
        if not row:
            _textbox(slide, left + Inches(0.25), Inches(3.5), Inches(5.4), Inches(0.4),
                     "Nothing to show yet.", size=13, color=MUTED)
            return
        y = Inches(2.4)
        for idx, item in enumerate(row, start=1):
            _add_round_rect(slide, left + Inches(0.18), y, Inches(5.5), Inches(0.72), BRAND_BLUE_SOFT)
            try:
                thumb = compose_element_thumb(item, cache=image_cache, size=180)
                _add_picture_bytes(
                    slide,
                    thumb,
                    left + Inches(0.28),
                    y + Inches(0.08),
                    width=Inches(0.56),
                    height=Inches(0.56),
                )
            except Exception:
                pass
            _textbox(
                slide,
                left + Inches(1.0),
                y + Inches(0.1),
                Inches(4.4),
                Inches(0.28),
                f"{idx}. {_safe(item.get('name'), 40)}",
                size=12,
                bold=True,
                color=NEAR_BLACK,
            )
            _textbox(
                slide,
                left + Inches(1.0),
                y + Inches(0.38),
                Inches(4.4),
                Inches(0.24),
                f"{_safe(item.get('category'), 28)}  ·  {_fmt_num(item.get('value'))}",
                size=10,
                color=MUTED,
            )
            y += Inches(0.78)

    def use_avoid():
        slide = _blank(prs)
        _slide_chrome(slide, "Use / avoid", "What to lean into and what to drop", page=13, total=15)
        _use_avoid_column(slide, use_items, left=Inches(0.55), title="Use", accent=ACCENT_GREEN)
        _use_avoid_column(slide, avoid_items, left=Inches(6.75), title="Avoid", accent=ACCENT_RED)

    def behaviour():
        slide = _blank(prs)
        _slide_chrome(slide, "Response behaviour", "Timing and engagement signals", page=14, total=15)
        _kpi_card(slide, Inches(0.7), Inches(1.8), Inches(5.7), Inches(1.4),
                  "Average response time", f"{_fmt_num(avg_rt)}s")
        fatigue_text = fatigue_result.get("answer_text") or "Fatigue summary unavailable."
        rt_text = rt_result.get("answer_text") or "Response time summary unavailable."
        _add_round_rect(slide, Inches(0.7), Inches(3.5), Inches(11.8), Inches(2.8), WHITE)
        _textbox(slide, Inches(1.0), Inches(3.7), Inches(11.2), Inches(0.35),
                 "Behaviour notes", size=15, bold=True, color=BRAND_BLUE)
        lines = [
            (_safe(rt_text, 220), {"size": 13, "space_after": 12, "color": NEAR_BLACK}),
            (_safe(fatigue_text, 220), {"size": 13, "space_after": 0, "color": SLATE}),
        ]
        _multi_lines(slide, Inches(1.0), Inches(4.2), Inches(11.2), Inches(1.8), lines)

    def closing():
        slide = _blank(prs)
        _add_rect(slide, 0, 0, SLIDE_W, SLIDE_H, BRAND_BLUE)
        _brand_wordmark_on_dark = slide.shapes.add_textbox(Inches(1.0), Inches(2.0), Inches(10), Inches(0.7))
        tf = _brand_wordmark_on_dark.text_frame
        p = tf.paragraphs[0]
        r1 = p.add_run()
        r1.text = "Mind"
        _set_run_font(r1, size=36, bold=True, color=WHITE)
        r2 = p.add_run()
        r2.text = "Surve"
        _set_run_font(r2, size=36, bold=True, color=BRAND_BLUE_MID)
        _textbox(slide, Inches(1.0), Inches(2.9), Inches(11), Inches(0.5),
                 "Recommended next steps", size=20, bold=True, color=WHITE)
        next_steps = [
            "1. Review the top design in the Design Configurator",
            "2. Stress-test winning elements across key segments",
            "3. Drop or redesign the weakest elements before launch",
            "4. Ask MindSurve AI for a deeper dive on any slide",
        ]
        y = Inches(3.6)
        for step in next_steps:
            _textbox(slide, Inches(1.0), y, Inches(11), Inches(0.35), step, size=14, color=WHITE)
            y += Inches(0.4)
        _textbox(
            slide,
            Inches(1.0),
            SLIDE_H - Inches(0.9),
            Inches(11),
            Inches(0.35),
            f"{study_title}  ·  Generated {generated_at}",
            size=12,
            color=BRAND_BLUE_MID,
        )

    # Build slides in order (~15)
    cover()
    agenda()
    overview()
    respondent_profile()
    executive()
    elements_slide(
        "Top 5 elements",
        "Highest-lift building blocks",
        top_elements,
        page=6,
        accent=ACCENT_GREEN,
    )
    elements_slide(
        "Worst 5 elements",
        "Lowest-lift building blocks",
        worst_elements,
        page=7,
        accent=ACCENT_RED,
    )
    designs_slide(
        "Top 5 designs",
        "Best complete combinations — configurator preview",
        top_design_pack,
        page=8,
        accent=ACCENT_GREEN,
    )
    designs_slide(
        "Worst 5 designs",
        "Weakest complete combinations — configurator preview",
        worst_design_pack,
        page=9,
        accent=ACCENT_RED,
    )
    classification()
    classification_detail()
    segments()
    use_avoid()
    behaviour()
    closing()

    buf = io.BytesIO()
    prs.save(buf)
    payload = buf.getvalue()
    meta = {
        "slide_count": len(prs.slides),
        "study_title": study_title,
        "panelists": panelists,
        "responses": responses,
        "filename": _pptx_filename(study_title),
    }
    return payload, meta


def _pptx_filename(study_title: str) -> str:
    base = re.sub(r"[^\w\s\-]+", "", study_title or "study").strip()
    base = re.sub(r"\s+", "-", base)[:60] or "study"
    stamp = datetime.now(timezone.utc).strftime("%Y%m%d")
    return f"{base}-MindSurve-analytics-{stamp}.pptx"
